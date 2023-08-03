#%%
from pptx import Presentation
import pptx
from pptx.util import Inches
import os
from PIL import Image
from html2image import Html2Image

def html_to_png(html_list,html_to_png_path):
    hti = Html2Image()
    hti = Html2Image(output_path=html_to_png_path)
    for i in range(len(html_list)):
        html = html_list[i]
        hti.screenshot(url=html, save_as=f'image{i}.png')

def createPPT(images):
    prs = Presentation()
    for image_path in images:
        slide_layout = prs.slide_layouts[5]  # Use slide layout with title and content
        slide = prs.slides.add_slide(slide_layout)
        
        left = Inches(0.5)
        top = pptx.util.Inches(1)
        width=pptx.util.Inches(9)
        height=pptx.util.Inches(5)
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        # add_picture(image_file, left, top, width=None, height=None)

    return prs

def get_image_paths(folder_path):
    print("program start")
    ret_arr = []
    image_name = os.listdir(folder_path)
    for name in image_name:
        if name[-1] == "g":
            image_path = folder_path+ f'\{name}' 
            print(image_path)
            ret_arr.append(image_path)
    return ret_arr

def preprocess(folder_path,path_of_resized_images):
    image_name = os.listdir(folder_path)
    for name in image_name:
        if name[-1] == "g":
            image_path = folder_path+ f'\{name}' 
            print(image_path)
            img = Image.open(image_path)
            # img.show()
            # img = img.resize((int(1024//1.6),int(700//1.6)))
            img = img.save(path_of_resized_images+ f'\{name}')

if __name__ == "__main__":
    html_list = [
    'https://i.insider.com/6461edf05cd814001883ccf0','https://i.insider.com/6461edf75cd814001883ccff','https://i.insider.com/6461edf25cd814001883ccf3','https://i.insider.com/6461edf85cd814001883cd05','https://i.insider.com/6461edf85cd814001883cd02','https://i.insider.com/6461edf25cd814001883ccf6','https://i.insider.com/6461edf711971c00188f4607','https://i.insider.com/6461edf311971c00188f45fe','https://i.insider.com/6461edf35cd814001883ccf9','https://i.insider.com/6461edf65cd814001883ccfc','https://i.insider.com/6461edf511971c00188f4604','https://i.insider.com/6461edf411971c00188f4601',
    'https://i.insider.com/644642c9101b77001857a500', 'https://i.insider.com/644642cc120cd7001865d1f1', 'https://i.insider.com/644642cb101b77001857a50f', 'https://i.insider.com/644642c8101b77001857a4fd', 'https://i.insider.com/644642ca120cd7001865d1df', 'https://i.insider.com/644642c9120cd7001865d1dc', 'https://i.insider.com/644642cb120cd7001865d1e3', 'https://i.insider.com/644642c9101b77001857a503', 'https://i.insider.com/644642cb120cd7001865d1e6', 'https://i.insider.com/644642cb120cd7001865d1e8', 'https://i.insider.com/644642cc120cd7001865d1ee', 'https://i.insider.com/644642ca101b77001857a506', 'https://i.insider.com/644642ca101b77001857a509', 'https://i.insider.com/644642cc120cd7001865d1eb', 'https://i.insider.com/644642ca101b77001857a50c', 'https://i.insider.com/644642cc101b77001857a512',
    'https://i.insider.com/6447e5636001ce00190642db', 'https://i.insider.com/6448cdd212a82b00192f9357', 'https://i.insider.com/6448ce4512a82b00192f936c', 'https://i.insider.com/6448ce3f12a82b00192f9360', 'https://i.insider.com/6448ce4112a82b00192f9363', 'https://i.insider.com/6448ce3690523900196fedad', 'https://i.insider.com/6448ce4290523900196fedb6', 'https://i.insider.com/6448ce3890523900196fedb0', 'https://i.insider.com/6448ce4290523900196fedb9', 'https://i.insider.com/6448ce4290523900196fedb4', 'https://i.insider.com/6448ce4312a82b00192f9366', 'https://i.insider.com/6448ce4390523900196fedbc', 'https://i.insider.com/6448ce4512a82b00192f9369', 'https://i.insider.com/6448ce4690523900196fedbf',
    'https://i.insider.com/644684b9101b77001857ab08', 'https://i.insider.com/644684b9120cd7001865d7a3', 'https://i.insider.com/644684b9101b77001857ab0b', 'https://i.insider.com/644684ba120cd7001865d7aa', 'https://i.insider.com/644684ba120cd7001865d7ad', 'https://i.insider.com/644684ba101b77001857ab12', 'https://i.insider.com/644684ba101b77001857ab14', 'https://i.insider.com/644684ba101b77001857ab0f', 'https://i.insider.com/644684ba120cd7001865d7a8', 'https://i.insider.com/644684ba120cd7001865d7af',
    'https://i.insider.com/6422f0af8c966b0019499d30', 'https://i.insider.com/6422f0af8c966b0019499d3a', 'https://i.insider.com/6422f0b28c966b0019499d55', 'https://i.insider.com/6422f0ad8c966b0019499d28', 'https://i.insider.com/6422f0af8c966b0019499d34', 'https://i.insider.com/6422f0b18c966b0019499d3f', 'https://i.insider.com/6422f0b18c966b0019499d42', 'https://i.insider.com/6422f0b28c966b0019499d59', 'https://i.insider.com/6422f0b28c966b0019499d53', 'https://i.insider.com/6422f0b18c966b0019499d46', 'https://i.insider.com/6422f0b28c966b0019499d51', 'https://i.insider.com/6422f0b28c966b0019499d61', 'https://i.insider.com/6422f0b28c966b0019499d5f', 'https://i.insider.com/6422f0b28c966b0019499d5d', 'https://i.insider.com/6422f0b18c966b0019499d4e', 'https://i.insider.com/6422f0af8c966b0019499d37', 'https://i.insider.com/6422f0b28c966b0019499d5b', 'https://i.insider.com/6422f0b28c966b0019499d57', 'https://i.insider.com/6422f0b18c966b0019499d44',
    'https://i.insider.com/62bc538f2224c30018142066', 'https://i.insider.com/62bc53912224c3001814206f', 'https://i.insider.com/62bc53922224c30018142079', 'https://i.insider.com/62bc53932224c30018142089', 'https://i.insider.com/62bc53922224c3001814207b', 'https://i.insider.com/62bc53922224c3001814207d', 'https://i.insider.com/62bc53922224c30018142083', 'https://i.insider.com/62bc53932224c3001814208b', 'https://i.insider.com/62bc53912224c30018142076', 'https://i.insider.com/62bc53912224c30018142072', 'https://i.insider.com/62bc53922224c3001814207f', 'https://i.insider.com/62bc53922224c30018142085', 'https://i.insider.com/62bc53922224c30018142081', 'https://i.insider.com/62bc53922224c30018142087', 'https://i.insider.com/62bc53912224c30018142074',
    'https://i.insider.com/636e1b65f5877200181c6b6d', 'https://i.insider.com/636e1b66951bdc00182da1dd', 'https://i.insider.com/636e1b66f5877200181c6b72', 'https://i.insider.com/636e1b67f5877200181c6b7a', 'https://i.insider.com/636e1b67951bdc00182da1e6', 'https://i.insider.com/636e1b67951bdc00182da1e8', 'https://i.insider.com/636e1b67f5877200181c6b76', 'https://i.insider.com/636e1b66951bdc00182da1df', 'https://i.insider.com/636e1b67f5877200181c6b78', 'https://i.insider.com/636e1c1bf5877200181c6b8e', 'https://i.insider.com/636e1b67951bdc00182da1e1', 'https://i.insider.com/636e1b67f5877200181c6b7e', 'https://i.insider.com/636e1b67f5877200181c6b7c', 'https://i.insider.com/636e1b65951bdc00182da1d7',
    'https://i.insider.com/63038de0c2794d00197ef91b', 'https://i.insider.com/63038debc2794d00197ef933', 'https://i.insider.com/63038ecac2794d00197ef9a0', 'https://i.insider.com/630dcbde2d6c740018c3f859', 'https://i.insider.com/630dcbdf4692300019737a15', 'https://i.insider.com/63038debc2794d00197ef936', 'https://i.insider.com/63038ecac2794d00197ef9a3', 'https://i.insider.com/63038ddbc9a2da0018090691', 'https://i.insider.com/63038de1c2794d00197ef924', 'https://i.insider.com/63038debc2794d00197ef931', 'https://i.insider.com/63038ecac9a2da0018090703', 'https://i.insider.com/63038ddfc2794d00197ef919', 'https://i.insider.com/63038de1c2794d00197ef921', 'https://i.insider.com/630dcbdf4692300019737a17',
    'https://i.insider.com/646c7713fee1b300197a596b', 'https://i.insider.com/646c7710fee1b300197a5955', 'https://i.insider.com/646c7712fee1b300197a5963', 'https://i.insider.com/646c770efee1b300197a5951', 'https://i.insider.com/646c7713fee1b300197a596d', 'https://i.insider.com/646c770dfee1b300197a5946', 'https://i.insider.com/646c7710fee1b300197a5959', 'https://i.insider.com/646c770bfee1b300197a593e', 'https://i.insider.com/646c770efee1b300197a594e', 'https://i.insider.com/646c7713fee1b300197a5970', 'https://i.insider.com/646c7713fee1b300197a5972', 'https://i.insider.com/646c7712fee1b300197a5967', 'https://i.insider.com/646c7714fee1b300197a5974', 'https://i.insider.com/646c7710fee1b300197a5957',
    'https://i.insider.com/64be4c6cbea34400195fb04f', 'https://i.insider.com/64be4c6dbea34400195fb058', 'https://i.insider.com/64be4c6eed4f46001961ada8', 'https://i.insider.com/64be4c6dbea34400195fb056', 'https://i.insider.com/64be4c6ded4f46001961ada4', 'https://i.insider.com/64be4c6eed4f46001961adad', 'https://i.insider.com/64be4c6eed4f46001961ada6', 'https://i.insider.com/64be4c6dbea34400195fb05b', 'https://i.insider.com/64be4c6ced4f46001961ad98',
    'https://i.insider.com/648899e8d4e551001974dc6a', 'https://i.insider.com/648899ead4e551001974dc78', 'https://i.insider.com/648899e8d4e551001974dc6d', 'https://i.insider.com/648899e9d4e551001974dc71', 'https://i.insider.com/648899e8d4e551001974dc68', 'https://i.insider.com/648899ead4e551001974dc7a', 'https://i.insider.com/648899ebd4e551001974dc85', 'https://i.insider.com/648899ebd4e551001974dc83', 'https://i.insider.com/648899ebd4e551001974dc7f', 'https://i.insider.com/648899ebd4e551001974dc81', 'https://i.insider.com/648899ead4e551001974dc76',
    'https://i.insider.com/642c07c9d335200018dd8470', 'https://i.insider.com/642c07cafcb86b0018030ac9', 'https://i.insider.com/642c07cbd335200018dd847b', 'https://i.insider.com/642c07cafcb86b0018030ace', 'https://i.insider.com/642c07cbd335200018dd847e', 'https://i.insider.com/642c07cbfcb86b0018030ad6', 'https://i.insider.com/642c07cad335200018dd8477', 'https://i.insider.com/642c07cbfcb86b0018030ad4', 'https://i.insider.com/642c07cbfcb86b0018030ad1', 'https://i.insider.com/642c07cad335200018dd8474', 'https://i.insider.com/642c07cad335200018dd8479', 'https://i.insider.com/642c07cbfcb86b0018030ad8', 'https://i.insider.com/642c07cafcb86b0018030acc',
    'https://i.insider.com/62e13e2ed8e3a400192ccdfc', 'https://i.insider.com/62e13e31d8e3a400192cce08', 'https://i.insider.com/62e13e30d8e3a400192cce02', 'https://i.insider.com/62e13e31d8e3a400192cce06', 'https://i.insider.com/62e13e33d8e3a400192cce11', 'https://i.insider.com/62e13e34d8e3a400192cce21', 'https://i.insider.com/62e13e34d8e3a400192cce23', 'https://i.insider.com/62e13e33d8e3a400192cce16', 'https://i.insider.com/62e13e33d8e3a400192cce0d', 'https://i.insider.com/62e13e33d8e3a400192cce14', 'https://i.insider.com/62e13e34d8e3a400192cce1d', 'https://i.insider.com/62e13e34d8e3a400192cce25', 'https://i.insider.com/62e13e34d8e3a400192cce1a', 'https://i.insider.com/62e13e34d8e3a400192cce1f', 'https://i.insider.com/62e13e34d8e3a400192cce27',
    'https://i.insider.com/64bfac56bea344001960dadc', 'https://i.insider.com/64bfac58bea344001960daf2', 'https://i.insider.com/64bfac58bea344001960dae9', 'https://i.insider.com/64bfac58bea344001960daf0', 'https://i.insider.com/64bfac58bea344001960daee', 'https://i.insider.com/64bfac58bea344001960daeb', 'https://i.insider.com/64bfac56bea344001960dadf', 'https://i.insider.com/64bfac58bea344001960daf4', 'https://i.insider.com/64bfac58bea344001960dae7',
    'https://i.insider.com/63c96138eee94d001a78eeea', 'https://i.insider.com/63c9613b2a1e600018b8b262', 'https://i.insider.com/63c9613aeee94d001a78eef2', 'https://i.insider.com/63c9613b2a1e600018b8b264', 'https://i.insider.com/63c9613beee94d001a78eefb', 'https://i.insider.com/63c9613b2a1e600018b8b26c', 'https://i.insider.com/63c9613beee94d001a78eef9', 'https://i.insider.com/63c9613b2a1e600018b8b266', 'https://i.insider.com/63c9613beee94d001a78eefd', 'https://i.insider.com/63c9613b2a1e600018b8b26a', 'https://i.insider.com/63c9613aeee94d001a78eef7', 'https://i.insider.com/63c9613b2a1e600018b8b268', 'https://i.insider.com/63c9613beee94d001a78eeff', 'https://i.insider.com/63c9613a2a1e600018b8b25f',
    'https://i.insider.com/63bbd90123f3620019876c94', 'https://i.insider.com/63bbd90023f3620019876c8e', 'https://i.insider.com/63bbd90106706e0019a96cc8', 'https://i.insider.com/63bbd90323f3620019876ca4', 'https://i.insider.com/63bbd90323f3620019876c9e', 'https://i.insider.com/63bbd90323f3620019876ca0', 'https://i.insider.com/63bbd90306706e0019a96cd1', 'https://i.insider.com/63bbd90223f3620019876c9a', 'https://i.insider.com/63bbd90406706e0019a96cd9', 'https://i.insider.com/63bbd90306706e0019a96cd4', 'https://i.insider.com/63bbd90406706e0019a96cdc', 'https://i.insider.com/63bbd90323f3620019876ca6', 'https://i.insider.com/63bbd90206706e0019a96ccc', 'https://i.insider.com/63bbd90306706e0019a96cd7', 'https://i.insider.com/63bbd90123f3620019876c92', 'https://i.insider.com/63bbd90106706e0019a96cc6', 'https://i.insider.com/63bbd90406706e0019a96cde', 'https://i.insider.com/63bbd90206706e0019a96cca', 'https://i.insider.com/63bbd90323f3620019876ca2',
    'https://i.insider.com/628e11773dbbd20018b71da1', 'https://i.insider.com/628e11783dbbd20018b71da5', 'https://i.insider.com/628e11797be71f001936f9c2', 'https://i.insider.com/628e11797be71f001936f9c0', 'https://i.insider.com/628e11797be71f001936f9c4', 'https://i.insider.com/628e11797be71f001936f9c6', 'https://i.insider.com/628e11783dbbd20018b71da7',
    'https://i.insider.com/6422ae738c966b00194995c9', 'https://i.insider.com/6422ae7355bd920018e202ad', 'https://i.insider.com/6422ae738c966b00194995cd', 'https://i.insider.com/6422ae798c966b00194995e0', 'https://i.insider.com/6422ae798c966b00194995de', 'https://i.insider.com/6422ae878c966b0019499644', 'https://i.insider.com/6422ae798c966b00194995dc', 'https://i.insider.com/6422ae7f8c966b0019499612', 'https://i.insider.com/6422ae7e8c966b0019499606', 'https://i.insider.com/6422ae898c966b0019499658', 'https://i.insider.com/6422ae7d8c966b0019499600', 'https://i.insider.com/6422ae7e8c966b0019499608', 'https://i.insider.com/6422ae7a8c966b00194995eb', 'https://i.insider.com/6422ae7e8c966b001949960a', 'https://i.insider.com/6422ae898c966b0019499654', 'https://i.insider.com/6422ae7c8c966b00194995f0', 'https://i.insider.com/6422ae858c966b0019499632', 'https://i.insider.com/6422ae828c966b001949961b', 'https://i.insider.com/6422ae7c8c966b00194995ed', 'https://i.insider.com/6422ae778c966b00194995d8', 'https://i.insider.com/6422ae7f8c966b0019499610', 'https://i.insider.com/6422ae828c966b0019499619', 'https://i.insider.com/6422ae828c966b0019499616', 'https://i.insider.com/6422ae7c8c966b00194995f2', 'https://i.insider.com/6422ae7c8c966b00194995fb', 'https://i.insider.com/6422ae7f8c966b0019499614', 'https://i.insider.com/6422ae898c966b001949964e', 'https://i.insider.com/6422ae7c8c966b00194995f5', 'https://i.insider.com/6422ae7a8c966b00194995e9', 'https://i.insider.com/6422ae848c966b001949962c', 'https://i.insider.com/6422ae828c966b001949961f', 'https://i.insider.com/6422ae868c966b001949963c', 'https://i.insider.com/6422ae7c8c966b00194995f9', 'https://i.insider.com/6422ae7e8c966b0019499604', 'https://i.insider.com/6422ae898c966b0019499650', 'https://i.insider.com/6422ae878c966b001949964c', 'https://i.insider.com/6422ae7d8c966b0019499602', 'https://i.insider.com/6422ae898c966b0019499656', 'https://i.insider.com/6422ae878c966b0019499642', 'https://i.insider.com/6422ae878c966b001949964a', 'https://i.insider.com/6422ae878c966b0019499648', 'https://i.insider.com/6422ae898c966b001949965a', 'https://i.insider.com/6422ae868c966b001949963e', 'https://i.insider.com/6422ae878c966b0019499640', 'https://i.insider.com/6422ae878c966b0019499646', 'https://i.insider.com/6422ae838c966b001949962a', 'https://i.insider.com/6422ae868c966b0019499638', 'https://i.insider.com/6422ae868c966b001949963a', 'https://i.insider.com/6422ae898c966b0019499652', 'https://i.insider.com/6422ae838c966b0019499628',
    'https://i.insider.com/63d15ac9def14100198e3b34', 'https://i.insider.com/63d15abf40fbb600192b6564', 'https://i.insider.com/63d15ac8def14100198e3b2c', 'https://i.insider.com/63d15ac2def14100198e3b0d', 'https://i.insider.com/63d15ac7def14100198e3b26', 'https://i.insider.com/63d15ac9def14100198e3b36', 'https://i.insider.com/63d15ac5def14100198e3b18', 'https://i.insider.com/63d15ac6def14100198e3b20', 'https://i.insider.com/63d15ac7def14100198e3b23', 'https://i.insider.com/63d15ac4def14100198e3b13', 'https://i.insider.com/63d15ac5def14100198e3b1a', 'https://i.insider.com/63d15ac0def14100198e3b08', 'https://i.insider.com/63d15ac5def14100198e3b1c', 'https://i.insider.com/63d15ac2def14100198e3b0f', 'https://i.insider.com/63d15ac8def14100198e3b2a', 'https://i.insider.com/63d15ac9def14100198e3b38', 'https://i.insider.com/63d15abf40fbb600192b6561',
    'https://i.insider.com/64a424564cc8540019cb5ae3', 'https://i.insider.com/64a424594cc8540019cb5b04', 'https://i.insider.com/64a424564cc8540019cb5ae6', 'https://i.insider.com/64a424594cc8540019cb5b00', 'https://i.insider.com/64a424584cc8540019cb5af2', 'https://i.insider.com/64a424594cc8540019cb5af9', 'https://i.insider.com/64a4245a4cc8540019cb5b0a', 'https://i.insider.com/64a424594cc8540019cb5afe', 'https://i.insider.com/64a424594cc8540019cb5afc', 'https://i.insider.com/64a4245a4cc8540019cb5b06', 'https://i.insider.com/64a4245a4cc8540019cb5b08', 'https://i.insider.com/64a424594cc8540019cb5af5', 'https://i.insider.com/64a424594cc8540019cb5b02', 'https://i.insider.com/64a424594cc8540019cb5af7',
    'https://i.insider.com/636e6509951bdc00182daa7b', 'https://i.insider.com/636e6510951bdc00182daaa0', 'https://i.insider.com/636e650f951bdc00182daa98', 'https://i.insider.com/636e6510951bdc00182daaaa', 'https://i.insider.com/636e650f951bdc00182daa96', 'https://i.insider.com/636e6510951bdc00182daa9e', 'https://i.insider.com/636e6510951bdc00182daaa4', 'https://i.insider.com/636e6510951bdc00182daaa2', 'https://i.insider.com/636e6510951bdc00182daaad', 'https://i.insider.com/636e6510951bdc00182daaa8', 'https://i.insider.com/636e650e951bdc00182daa91', 'https://i.insider.com/636e6510951bdc00182daaaf', 'https://i.insider.com/636e650f951bdc00182daa9a', 'https://i.insider.com/636e650f951bdc00182daa94', 'https://i.insider.com/636e6510951bdc00182daaa6', 'https://i.insider.com/636e650f951bdc00182daa9c'
    ]


    html_to_png_path = r'C:\Users\stanl\Desktop\html_to_png'
    path_of_resized_images = r'C:\Users\stanl\Desktop\resized_images'

    html_to_png(html_list,html_to_png_path)

    # preprocess(html_to_png_path,path_of_resized_images)
    
    image_path_arr = get_image_paths(html_to_png_path)

    prs = createPPT(image_path_arr)

    prs.save("output_presentation.pptx")
    print("PowerPoint presentation saved successfully!")

    

#%%
